import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

# ==============================================================================
# EOJOB BRAND COLOR SYSTEM (이어잡 공식 컬러 팔레트)
# ==============================================================================
COLOR_PRIMARY       = RGBColor(23, 63, 58)       # #173F3A (Deep Evergreen - 깊고 신뢰감 있는 숲색)
COLOR_PRIMARY_DARK  = RGBColor(15, 42, 39)       # #0F2A27 (Dark Evergreen)
COLOR_PRIMARY_LIGHT = RGBColor(36, 90, 83)       # #245A53 (Medium Teal)
COLOR_ACCENT        = RGBColor(240, 107, 79)     # #F06B4F (Accessible Coral - 역동적인 산호색)
COLOR_ACCENT_DARK   = RGBColor(184, 71, 52)      # #B84734 (Dark Coral)
COLOR_ACCENT_BG     = RGBColor(253, 240, 237)    # #FDF0ED (Subtle Coral Tint)
COLOR_IVORY         = RGBColor(247, 243, 234)    # #F7F3EA (Warm Ivory - 따뜻하고 세련된 배경)
COLOR_WHITE         = RGBColor(255, 255, 255)    # #FFFFFF (Pure White)
COLOR_MINT          = RGBColor(221, 235, 231)    # #DDEBE7 (Soft Mint)
COLOR_MINT_BORDER   = RGBColor(187, 213, 206)    # #BBD5CE
COLOR_BORDER        = RGBColor(224, 217, 200)    # #E0D9C8 (Warm Border)
COLOR_BORDER_LIGHT  = RGBColor(238, 233, 222)    # #EEE9DE
COLOR_INK           = RGBColor(23, 33, 43)       # #17212B (Deep Slate Ink)
COLOR_INK_MUTED     = RGBColor(83, 96, 110)      # #53606E (Muted Ink)
COLOR_INK_SUBTLE    = RGBColor(138, 150, 163)    # #8A96A3 (Subtle Gray)

FONT_FAMILY = "Pretendard"

# ==============================================================================
# OpenXML FONT HELPER (한글/영문 모두 100% Pretendard 강제 주입)
# ==============================================================================
def apply_pretendard_to_run(run):
    run.font.name = FONT_FAMILY
    rPr = run._r.get_or_add_rPr()
    for tag in ['ea', 'latin', 'cs']:
        el = rPr.find(qn(f'a:{tag}'))
        if el is None:
            el = parse_xml(f'<a:{tag} {nsdecls("a")} typeface="{FONT_FAMILY}"/>')
            rPr.append(el)
        else:
            el.set('typeface', FONT_FAMILY)

def apply_pretendard_to_paragraph(paragraph):
    for run in paragraph.runs:
        apply_pretendard_to_run(run)
    pPr = paragraph._p.get_or_add_pPr()
    defRPr = pPr.find(qn('a:defRPr'))
    if defRPr is None:
        defRPr = parse_xml(f'<a:defRPr {nsdecls("a")}><a:latin typeface="{FONT_FAMILY}"/><a:ea typeface="{FONT_FAMILY}"/><a:cs typeface="{FONT_FAMILY}"/></a:defRPr>')
        pPr.append(defRPr)
    else:
        for tag in ['latin', 'ea', 'cs']:
            f_el = defRPr.find(qn(f'a:{tag}'))
            if f_el is None:
                f_el = parse_xml(f'<a:{tag} {nsdecls("a")} typeface="{FONT_FAMILY}"/>')
                defRPr.append(f_el)
            else:
                f_el.set('typeface', FONT_FAMILY)

def enforce_pretendard_everywhere(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    apply_pretendard_to_paragraph(p)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            apply_pretendard_to_paragraph(p)

# ==============================================================================
# DESIGN COMPONENT BUILDERS (세련된 디자인 요소 헬퍼)
# ==============================================================================
def set_bg_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_header(slide, title, category_num, category_name, is_dark=False):
    # Category Eyebrow Pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(2.2), Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT if is_dark else COLOR_MINT
    pill.line.color.rgb = COLOR_MINT if is_dark else COLOR_MINT_BORDER
    pill.line.width = Pt(1)
    ptf = pill.text_frame
    ptf.margin_top = Inches(0.04)
    pp = ptf.paragraphs[0]
    pp.text = f"{category_num}  |  {category_name.upper()}"
    pp.font.name = FONT_FAMILY
    pp.font.size = Pt(9.5)
    pp.font.bold = True
    pp.font.color.rgb = COLOR_WHITE if is_dark else COLOR_PRIMARY
    pp.alignment = PP_ALIGN.CENTER

    # Slide Title
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(10.5), Inches(0.6))
    ttf = t_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.name = FONT_FAMILY
    tp.font.size = Pt(21)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_WHITE if is_dark else COLOR_PRIMARY

    # Header Right Mini Tag
    tag_box = slide.shapes.add_textbox(Inches(9.2), Inches(0.5), Inches(3.3), Inches(0.35))
    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = "EOJOB B2B TALENT SOLUTION"
    tag_p.font.name = FONT_FAMILY
    tag_p.font.size = Pt(9.5)
    tag_p.font.bold = True
    tag_p.font.color.rgb = COLOR_ACCENT if is_dark else COLOR_INK_MUTED
    tag_p.alignment = PP_ALIGN.RIGHT

    # Elegant divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT if is_dark else COLOR_BORDER
    line.line.fill.background()

# ==============================================================================
# MAIN PRESENTATION BUILDER
# ==============================================================================
def create_high_end_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    assets_dir = os.path.abspath("docs/slides/assets")
    logo_path = os.path.join(assets_dir, "logo_text.png")
    img_senior = os.path.join(assets_dir, "eojob_senior_leader.jpg")
    img_hr = os.path.join(assets_dir, "eojob_hr_concept.jpg")
    img_ai_card = os.path.join(assets_dir, "eojob_ai_card_mockup.jpg")
    img_success = os.path.join(assets_dir, "eojob_success_case.jpg")

    # --------------------------------------------------------------------------
    # SLIDE 1: HERO 표지 (High-End Dark Evergreen Theme)
    # --------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_bg_color(s1, COLOR_PRIMARY)

    # Decorative background geometric accents (입체감 부여)
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()

    # White Logo in Cover
    if os.path.exists(logo_path):
        # 로고 박스
        s1.shapes.add_picture(logo_path, Inches(1.0), Inches(1.0), height=Inches(0.46))

    # Category Pill
    c_pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.75), Inches(3.6), Inches(0.36))
    c_pill.fill.solid()
    c_pill.fill.fore_color.rgb = COLOR_PRIMARY_LIGHT
    c_pill.line.color.rgb = COLOR_ACCENT
    c_pill.line.width = Pt(1.5)
    cp_tf = c_pill.text_frame
    cp_p = cp_tf.paragraphs[0]
    cp_p.text = "FOR HR EXECUTIVES & RECRUITERS"
    cp_p.font.name = FONT_FAMILY
    cp_p.font.size = Pt(10)
    cp_p.font.bold = True
    cp_p.font.color.rgb = COLOR_WHITE
    cp_p.alignment = PP_ALIGN.CENTER

    # Hero Massive Typography
    t1_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(6.8), Inches(2.4))
    t1_tf = t1_box.text_frame
    t1_tf.word_wrap = True

    p1 = t1_tf.paragraphs[0]
    p1.text = "실무 검증된 시니어 전문가,"
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE

    p2 = t1_tf.add_paragraph()
    p2.text = "정부 지원금으로 연 720만원\n인건비를 절감하세요"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT

    # Cover Description
    d_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(6.5), Inches(1.0))
    d_tf = d_box.text_frame
    d_tf.word_wrap = True
    dp = d_tf.paragraphs[0]
    dp.text = "이어잡(EOJOB)은 AI 심층 역량 분석을 통해 당면 과제를 즉시 해결할 인재만 선별하고, 고용노동부 고용촉진장려금 환급을 원스톱으로 연결하는 기업 전용 인재 솔루션입니다."
    dp.font.name = FONT_FAMILY
    dp.font.size = Pt(12.5)
    dp.font.color.rgb = COLOR_MINT

    # 3 Stat Badges in Cover
    b_x = Inches(1.0)
    badges = [
        ("서류 합격률", "2.4배 ↑"),
        ("1인당 지원금", "연 720만원"),
        ("채용 리스크", "14일 무상 보증"),
    ]
    for b_lbl, b_val in badges:
        b_shape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, b_x, Inches(5.8), Inches(2.05), Inches(0.95))
        b_shape.fill.solid()
        b_shape.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        b_shape.line.color.rgb = COLOR_PRIMARY_LIGHT
        b_shape.line.width = Pt(1)
        btf = b_shape.text_frame
        btf.margin_top = Inches(0.12)
        bp1 = btf.paragraphs[0]
        bp1.text = b_lbl
        bp1.font.name = FONT_FAMILY
        bp1.font.size = Pt(9.5)
        bp1.font.color.rgb = COLOR_INK_SUBTLE
        bp1.alignment = PP_ALIGN.CENTER
        bp2 = btf.add_paragraph()
        bp2.text = b_val
        bp2.font.name = FONT_FAMILY
        bp2.font.size = Pt(14)
        bp2.font.bold = True
        bp2.font.color.rgb = COLOR_ACCENT if "720" in b_val else COLOR_WHITE
        bp2.alignment = PP_ALIGN.CENTER
        b_x += Inches(2.25)

    # Right Hero Portrait Image Card
    if os.path.exists(img_senior):
        # Card Background Frame
        f_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.1), Inches(4.4), Inches(5.6))
        f_card.fill.solid()
        f_card.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        f_card.line.color.rgb = COLOR_PRIMARY_LIGHT
        f_card.line.width = Pt(1.5)

        s1.shapes.add_picture(img_senior, Inches(8.25), Inches(1.25), width=Inches(4.1), height=Inches(4.5))

        # Floating Trust Badge over photo
        t_pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(5.2), Inches(3.8), Inches(0.65))
        t_pill.fill.solid()
        t_pill.fill.fore_color.rgb = COLOR_WHITE
        t_pill.line.color.rgb = COLOR_ACCENT
        t_pill.line.width = Pt(1.5)
        tptf = t_pill.text_frame
        tptf.margin_top = Inches(0.08)
        tpp1 = tptf.paragraphs[0]
        tpp1.text = "🎖️ 20년+ 대기업/공공 출신 검증된 인재 풀"
        tpp1.font.name = FONT_FAMILY
        tpp1.font.size = Pt(10.5)
        tpp1.font.bold = True
        tpp1.font.color.rgb = COLOR_PRIMARY
        tpp1.alignment = PP_ALIGN.CENTER
        tpp2 = tptf.add_paragraph()
        tpp2.text = "제조품질 · 클라우드 IT · 국책 R&D 즉시 투입"
        tpp2.font.name = FONT_FAMILY
        tpp2.font.size = Pt(9)
        tpp2.font.color.rgb = COLOR_INK_MUTED
        tpp2.alignment = PP_ALIGN.CENTER

    # --------------------------------------------------------------------------
    # SLIDE 2: 채용 위기 데이터 (High-Impact Stat Cards)
    # --------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_bg_color(s2, COLOR_IVORY)
    create_header(s2, "경력직 1명 채용에 850만원, 그런데 3개월 내 조기 퇴사율 27%", "01", "MARKET REALITY")

    # 3 Giant KPI Cards
    kpi_x = Inches(0.8)
    kpis = [
        ("850 만원", "평균 경력 채용 소요 비용", "서칭, 공고비, 헤드헌팅 수수료(연봉의 15~25%) 및 서류 검토에 드는 기업 실질 지출액", COLOR_WHITE, COLOR_ACCENT),
        ("27.4 %", "입사 3개월 내 조기 퇴사율", "화려한 출신 직급만 보고 채용했으나, 실제 프로젝트 과제 미스매치로 인한 퇴사 손실", COLOR_WHITE, RGBColor(184, 71, 52)),
        ("720 만원", "1인당 미수령 고용촉진장려금", "복잡한 서류 절차와 대상자 파악의 어려움으로 매년 기업들이 놓치는 정부 환급금", COLOR_WHITE, COLOR_PRIMARY),
    ]

    for num_str, title_str, desc_str, bg_c, num_c in kpis:
        k_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, kpi_x, Inches(1.8), Inches(3.7), Inches(2.3))
        k_card.fill.solid()
        k_card.fill.fore_color.rgb = bg_c
        k_card.line.color.rgb = COLOR_BORDER
        k_card.line.width = Pt(1.5)

        ktf = k_card.text_frame
        ktf.word_wrap = True
        ktf.margin_left = Inches(0.25)
        ktf.margin_top = Inches(0.2)

        kp1 = ktf.paragraphs[0]
        kp1.text = num_str
        kp1.font.name = FONT_FAMILY
        kp1.font.size = Pt(28)
        kp1.font.bold = True
        kp1.font.color.rgb = num_c

        kp2 = ktf.add_paragraph()
        kp2.text = title_str
        kp2.font.name = FONT_FAMILY
        kp2.font.size = Pt(13)
        kp2.font.bold = True
        kp2.font.color.rgb = COLOR_PRIMARY

        kp3 = ktf.add_paragraph()
        kp3.text = "\n" + desc_str
        kp3.font.name = FONT_FAMILY
        kp3.font.size = Pt(10.5)
        kp3.font.color.rgb = COLOR_INK_MUTED

        kpi_x += Inches(4.0)

    # Bottom Split: Concept Image + Core Solution Callout
    bot_y = Inches(4.35)
    if os.path.exists(img_hr):
        s2.shapes.add_picture(img_hr, Inches(0.8), bot_y, width=Inches(4.2), height=Inches(2.55))

    sol_box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), bot_y, Inches(7.333), Inches(2.55))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = COLOR_PRIMARY
    sol_box.line.fill.background()

    stf = sol_box.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.35)
    stf.margin_top = Inches(0.25)

    sp1 = stf.paragraphs[0]
    sp1.text = "💡 이어잡(EOJOB)의 혁신적인 3대 해결책"
    sp1.font.name = FONT_FAMILY
    sp1.font.size = Pt(15)
    sp1.font.bold = True
    sp1.font.color.rgb = COLOR_WHITE

    sp2 = stf.add_paragraph()
    sp2.text = "\n1. 스펙 중심 채용 ➔ AI 음성 인터뷰 기반 '문제 해결 역량 검증 카드' 사전 제공"
    sp2.font.name = FONT_FAMILY
    sp2.font.size = Pt(11.5)
    sp2.font.color.rgb = COLOR_MINT

    sp3 = stf.add_paragraph()
    sp3.text = "2. 경직된 전일제 정규직 ➔ 주 2~3일 하이브리드 고문 / 파트타임 유연 계약 지원"
    sp3.font.name = FONT_FAMILY
    sp3.font.size = Pt(11.5)
    sp3.font.color.rgb = COLOR_MINT

    sp4 = stf.add_paragraph()
    sp4.text = "3. 번거로운 서류 작업 ➔ 채용 즉시 분기별 180만원(연 720만원) 자동 신청 리포트 발행"
    sp4.font.name = FONT_FAMILY
    sp4.font.size = Pt(11.5)
    sp4.font.color.rgb = COLOR_ACCENT

    # --------------------------------------------------------------------------
    # SLIDE 3: 비교 매트릭스 표 (High-End SaaS Feature Matrix)
    # --------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_bg_color(s3, COLOR_IVORY)
    create_header(s3, "기존 채용 방식 vs 이어잡(EOJOB) 솔루션 5대 핵심 지표 비교", "02", "COMPETITIVE ADVANTAGE")

    # Table Shape
    tbl_shape = s3.shapes.add_table(6, 4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
    tbl = tbl_shape.table

    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(3.0)
    tbl.columns[2].width = Inches(3.0)
    tbl.columns[3].width = Inches(3.533)

    headers = ["비교 지표", "일반 채용 포털", "전통 헤드헌팅 사", "이어잡 (EOJOB)  ★ 추천"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY if i == 3 else COLOR_MINT
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE if i == 3 else COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER

    data = [
        ("역량 검증 방식", "지원자의 일방적 텍스트 이력서", "단순 연차 및 지인 평판 조회", "AI 인터뷰 기반 '문제 해결 카드' 제공"),
        ("고용 형태 유연성", "정규직 위주 (경직성 높음)", "정규직 전용 (수수료 중심)", "정규직 + 주 2~3일 유연 자문 계약"),
        ("정부 지원금 연계", "기업이 직접 제도 파악 및 신청", "연계 지원 전혀 없음", "연 720만원(월 60만원) 자동 산출서"),
        ("채용 리스크 보장", "보장 없음 (환불/교체 불가)", "3개월 내 1회 한정 재추천", "14일 무상 교체 보증 (Risk-Zero)"),
        ("매칭 소요 시간", "서류 검토 평균 3주 이상 소요", "후보자 서칭 평균 4~6주", "AI 정밀 매칭 48시간 내 후보자 전달"),
    ]

    for r_idx, r_data in enumerate(data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            cell.fill.solid()
            # Highlight EOJOB column in distinct subtle tint
            cell.fill.fore_color.rgb = COLOR_ACCENT_BG if c_idx == 3 else COLOR_WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_FAMILY
            p.font.size = Pt(11)
            p.font.bold = (c_idx == 3 or c_idx == 0)
            if c_idx == 3:
                p.font.color.rgb = COLOR_ACCENT_DARK if r_idx in [3, 4, 5] else COLOR_PRIMARY
            else:
                p.font.color.rgb = COLOR_INK
            p.alignment = PP_ALIGN.CENTER

    # Bottom Impact Strip
    strip = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.48))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_PRIMARY
    strip.line.fill.background()
    stp = strip.text_frame.paragraphs[0]
    stp.text = "🎯 결론: 검증 시간은 80% 단축하고, 정부 지원금으로 채용 1인당 연 720만원 인건비를 즉각 보전합니다."
    stp.font.name = FONT_FAMILY
    stp.font.size = Pt(11.5)
    stp.font.bold = True
    stp.font.color.rgb = COLOR_WHITE
    stp.alignment = PP_ALIGN.CENTER

    # --------------------------------------------------------------------------
    # SLIDE 4: AI 경험 카드 & 인재 풀 쇼케이스 (High-Tech Visual Layout)
    # --------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_bg_color(s4, COLOR_IVORY)
    create_header(s4, "검토 시간 80% 단축, 실무 역량만 사전 추출하는 AI 경험 카드", "03", "AI TALENT CARD")

    # Left: AI Card 3D Mockup
    if os.path.exists(img_ai_card):
        s4.shapes.add_picture(img_ai_card, Inches(0.8), Inches(1.8), width=Inches(5.5), height=Inches(5.1))

    # Right: Talent Showcase Cards
    tx = Inches(6.6)
    ty = Inches(1.8)

    t_header_box = s4.shapes.add_textbox(tx, ty, Inches(5.9), Inches(0.6))
    th_tf = t_header_box.text_frame
    th_tf.word_wrap = True
    th_p = th_tf.paragraphs[0]
    th_p.text = "즉시 투입 가능한 상위 5% 시니어 전문가 네트워크"
    th_p.font.name = FONT_FAMILY
    th_p.font.size = Pt(14)
    th_p.font.bold = True
    th_p.font.color.rgb = COLOR_PRIMARY

    talents = [
        ("제조 품질 · 스마트공장 총괄 전무 (28년)", "완성차 1차 협력사 품질총괄 | ISO 9001 인증 고도화", "불량률 4.2% ➔ 0.8% 개선, 완성차 납기 클레임 0건 달성", "연 720만원 지원 대상", COLOR_ACCENT),
        ("클라우드 / 인프라 총괄 아키텍트 (22년)", "대기업 SI 및 테크 유니콘 인프라 리드 역임", "AWS/Azure 대규모 마이그레이션 & 인프라 유지비 35% 절감", "주 2~3일 유연 자문", COLOR_PRIMARY),
        ("정부 R&D 국책과제 기획 총괄 (19년)", "중기부 / 산자부 대형 국책과제 42억원 수주 총괄", "과제 기획서 작성, 대면 발표 평가 완벽 코칭", "프로젝트 단위 계약", COLOR_PRIMARY_LIGHT),
    ]

    ty += Inches(0.65)
    for name_str, sub_str, achieve_str, tag_str, tag_color in talents:
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, ty, Inches(5.9), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1.5)

        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.2)
        ctf.margin_top = Inches(0.12)

        p1 = ctf.paragraphs[0]
        p1.text = name_str
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY

        p2 = ctf.add_paragraph()
        p2.text = sub_str
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_INK_MUTED

        p3 = ctf.add_paragraph()
        p3.text = "성과: " + achieve_str
        p3.font.name = FONT_FAMILY
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_INK

        # Mini Badge on top right of card
        tag_shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx + Inches(4.3), ty + Inches(0.12), Inches(1.45), Inches(0.28))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = COLOR_MINT if tag_color != COLOR_ACCENT else COLOR_ACCENT_BG
        tag_shape.line.color.rgb = tag_color
        tag_shape.line.width = Pt(1)
        tt_p = tag_shape.text_frame.paragraphs[0]
        tt_p.text = tag_str
        tt_p.font.name = FONT_FAMILY
        tt_p.font.size = Pt(8.5)
        tt_p.font.bold = True
        tt_p.font.color.rgb = tag_color
        tt_p.alignment = PP_ALIGN.CENTER

        ty += Inches(1.45)

    # --------------------------------------------------------------------------
    # SLIDE 5: 고용촉진장려금 연 720만원 완벽 구조 (Financial ROI Dashboard)
    # --------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_bg_color(s5, COLOR_IVORY)
    create_header(s5, "고용촉진장려금 지원 구조 및 채용 인원별 현금 환급 시뮬레이션", "04", "FINANCIAL BENEFIT")

    # Left: Mechanism Card
    mech_card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.3), Inches(4.9))
    mech_card.fill.solid()
    mech_card.fill.fore_color.rgb = COLOR_PRIMARY
    mech_card.line.fill.background()

    mtf = mech_card.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.3)
    mtf.margin_top = Inches(0.3)

    mp1 = mtf.paragraphs[0]
    mp1.text = "고용노동부 고용촉진장려금 핵심 요약"
    mp1.font.name = FONT_FAMILY
    mp1.font.size = Pt(16)
    mp1.font.bold = True
    mp1.font.color.rgb = COLOR_WHITE

    details = [
        ("지원 대상", "국민취업지원제도 이수 인증 시니어 인재 채용 시"),
        ("지원 금액", "채용 1인당 월 60만원 (연간 최대 720만원)"),
        ("지급 주기", "3개월 단위로 분기별 180만원씩 회사 계좌로 직접 환급"),
        ("법적 근거", "고용보험법 시행령 제12조 (우선지원대상기업 기준)"),
    ]

    for d_title, d_val in details:
        p_t = mtf.add_paragraph()
        p_t.text = f"\n• {d_title}"
        p_t.font.name = FONT_FAMILY
        p_t.font.size = Pt(11.5)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_ACCENT

        p_v = mtf.add_paragraph()
        p_v.text = f"  {d_val}"
        p_v.font.name = FONT_FAMILY
        p_v.font.size = Pt(11)
        p_v.font.color.rgb = COLOR_MINT

    p_guide = mtf.add_paragraph()
    p_guide.text = "\n🛡️ 이어잡 플랫폼은 대상자 사전 선별부터 고용24 신청 서류 가이드를 원스톱으로 지원합니다."
    p_guide.font.name = FONT_FAMILY
    p_guide.font.size = Pt(11)
    p_guide.font.bold = True
    p_guide.font.color.rgb = COLOR_WHITE

    # Right: 4 Simulation Tiles
    sim_x = Inches(6.4)
    sim_y = Inches(1.8)
    sim_tiles = [
        ("1명 채용 시", "분기 180만원 환급", "연 720 만원", "신입 1인 인건비의 약 25% 보전 효과"),
        ("2명 채용 시", "분기 360만원 환급", "연 1,440 만원", "연간 1,440만원 순 현금 회사 통장 입금"),
        ("3명 채용 시", "분기 540만원 환급", "연 2,160 만원", "시니어 전문 자문단 구축 비용 대폭 절감"),
        ("5명 채용 시", "분기 900만원 환급", "연 3,600 만원", "1개 부서 인건비의 30% 이상 정부 보전"),
    ]

    for t_idx, (t_head, t_sub, t_amt, t_effect) in enumerate(sim_tiles):
        tile = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sim_x, sim_y, Inches(6.1), Inches(1.12))
        tile.fill.solid()
        tile.fill.fore_color.rgb = COLOR_WHITE
        tile.line.color.rgb = COLOR_BORDER
        tile.line.width = Pt(1.5)

        ttf = tile.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Inches(0.2)
        ttf.margin_top = Inches(0.12)

        tp1 = ttf.paragraphs[0]
        tp1.text = f"{t_head}  |  {t_sub}"
        tp1.font.name = FONT_FAMILY
        tp1.font.size = Pt(11)
        tp1.font.bold = True
        tp1.font.color.rgb = COLOR_PRIMARY

        tp2 = ttf.add_paragraph()
        tp2.text = t_effect
        tp2.font.name = FONT_FAMILY
        tp2.font.size = Pt(10)
        tp2.font.color.rgb = COLOR_INK_MUTED

        # Big Amount Badge on right of tile
        amt_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sim_x + Inches(4.3), sim_y + Inches(0.2), Inches(1.6), Inches(0.7))
        amt_box.fill.solid()
        amt_box.fill.fore_color.rgb = COLOR_ACCENT_BG
        amt_box.line.color.rgb = COLOR_ACCENT
        amt_box.line.width = Pt(1.5)
        ap = amt_box.text_frame.paragraphs[0]
        ap.text = t_amt
        ap.font.name = FONT_FAMILY
        ap.font.size = Pt(14)
        ap.font.bold = True
        ap.font.color.rgb = COLOR_ACCENT_DARK
        ap.alignment = PP_ALIGN.CENTER

        sim_y += Inches(1.25)

    # --------------------------------------------------------------------------
    # SLIDE 6: 실제 도입 성공 사례 (Case Study Infographic)
    # --------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_bg_color(s6, COLOR_IVORY)
    create_header(s6, "스마트제조 S사 도입 성공 사례 : 불량률 72% 감축 & 1,440만원 환급", "05", "PROVEN CASE STUDY")

    # Left Narrative Flow (Problem ➔ Solution ➔ Outcome)
    cy = Inches(1.8)
    cases = [
        ("1. 직면 과제 (Challenge)", "신규 자동차 정밀 부품 양산 라인에서 원인 불명의 불량률(2.5%) 지속 발생\n• 1차 협력사 납기 지연 위기 및 연간 폐기 손실 3.2억원 발생 위험", COLOR_WHITE, RGBColor(254, 178, 178), COLOR_ACCENT_DARK),
        ("2. 이어잡 매칭 솔루션 (Solution)", "완성차 1차사 품질총괄 28년 경력 시니어 전문가 2인 주 3일 하이브리드 투입\n• 48시간 내 매칭 완료 및 공정별 통계적 공정관리(SPC) 표준화 전수", COLOR_WHITE, COLOR_MINT_BORDER, COLOR_PRIMARY),
        ("3. 달성 성과 (Outcome)", "불량률 2.5% ➔ 0.15% 급감 (불량 72% 이상 감축 달성)\n• 연간 폐기 비용 3.2억원 절감 + 고용노동부 장려금 연 1,440만원 전액 수령", COLOR_ACCENT_BG, COLOR_ACCENT, COLOR_ACCENT_DARK),
    ]

    for c_title, c_desc, bg_c, b_c, title_c in cases:
        c_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), cy, Inches(6.0), Inches(1.5))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = bg_c
        c_box.line.color.rgb = b_c
        c_box.line.width = Pt(1.5)

        ctf = c_box.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.25)
        ctf.margin_top = Inches(0.18)

        cp1 = ctf.paragraphs[0]
        cp1.text = c_title
        cp1.font.name = FONT_FAMILY
        cp1.font.size = Pt(12.5)
        cp1.font.bold = True
        cp1.font.color.rgb = title_c

        cp2 = ctf.add_paragraph()
        cp2.text = c_desc
        cp2.font.name = FONT_FAMILY
        cp2.font.size = Pt(10.5)
        cp2.font.color.rgb = COLOR_INK

        cy += Inches(1.68)

    # Right: Success Case Photo & Milestone Box
    if os.path.exists(img_success):
        s6.shapes.add_picture(img_success, Inches(7.1), Inches(1.8), width=Inches(5.4), height=Inches(3.3))

    m_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(5.25), Inches(5.4), Inches(1.45))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = COLOR_PRIMARY
    m_box.line.fill.background()

    mtf = m_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.25)
    mtf.margin_top = Inches(0.15)

    mp1 = mtf.paragraphs[0]
    mp1.text = "💬 스마트제조 S사 인사총괄 전무의 한마디"
    mp1.font.name = FONT_FAMILY
    mp1.font.size = Pt(11)
    mp1.font.bold = True
    mp1.font.color.rgb = COLOR_ACCENT

    mp2 = mtf.add_paragraph()
    mp2.text = "\"기존 헤드헌팅으로는 찾을 수 없었던 28년 경력의 베테랑을 주 3일 자문 형태로 합리적으로 채용했고, 지원금 1,440만원까지 지원받아 인건비 부담이 전혀 없었습니다.\""
    mp2.font.name = FONT_FAMILY
    mp2.font.size = Pt(10.5)
    mp2.font.color.rgb = COLOR_WHITE

    # --------------------------------------------------------------------------
    # SLIDE 7: 3단계 채용 절차 & 14일 안심 보증제 (Process & Guarantee)
    # --------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_bg_color(s7, COLOR_IVORY)
    create_header(s7, "단 3단계 채용 프로세스 및 인사담당자 100% 안심 보증제", "06", "PROCESS & GUARANTEE")

    # 3 Step Flow Horizontal Cards
    sx = Inches(0.8)
    steps = [
        ("01", "과제 / 포지션 등록", "1분 등록", "해결하고 싶은 기술 과제나 채용 요건을 등록하면, AI가 요구 역량을 자동 파싱합니다."),
        ("02", "AI 맞춤 추천 & 검토", "48시간 내 완료", "적합도 90점 이상 & 연 720만원 지원 대상 인재의 AI 경험 카드를 전달받습니다."),
        ("03", "계약 & 지원금 수령", "원스톱 지원", "정규직/유연 계약 체결 후 고용24 분기별 180만원 환급 절차를 밀착 지원합니다."),
    ]

    for num, title, badge, desc in steps:
        s_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(1.8), Inches(3.7), Inches(2.5))
        s_card.fill.solid()
        s_card.fill.fore_color.rgb = COLOR_WHITE
        s_card.line.color.rgb = COLOR_ACCENT if num == "02" else COLOR_BORDER
        s_card.line.width = Pt(2 if num == "02" else 1.5)

        stf = s_card.text_frame
        stf.word_wrap = True
        stf.margin_left = Inches(0.25)
        stf.margin_top = Inches(0.2)

        sp1 = stf.paragraphs[0]
        sp1.text = f"STEP {num}"
        sp1.font.name = FONT_FAMILY
        sp1.font.size = Pt(11)
        sp1.font.bold = True
        sp1.font.color.rgb = COLOR_ACCENT if num == "02" else COLOR_PRIMARY_LIGHT

        sp2 = stf.add_paragraph()
        sp2.text = title
        sp2.font.name = FONT_FAMILY
        sp2.font.size = Pt(14)
        sp2.font.bold = True
        sp2.font.color.rgb = COLOR_PRIMARY

        sp3 = stf.add_paragraph()
        sp3.text = f"[{badge}]"
        sp3.font.name = FONT_FAMILY
        sp3.font.size = Pt(9.5)
        sp3.font.bold = True
        sp3.font.color.rgb = COLOR_ACCENT_DARK

        sp4 = stf.add_paragraph()
        sp4.text = "\n" + desc
        sp4.font.name = FONT_FAMILY
        sp4.font.size = Pt(10.5)
        sp4.font.color.rgb = COLOR_INK_MUTED

        sx += Inches(4.0)

    # Risk-Zero Guarantee Shield Card (대형 안심 보증 배너)
    g_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(11.733), Inches(2.1))
    g_card.fill.solid()
    g_card.fill.fore_color.rgb = COLOR_PRIMARY
    g_card.line.color.rgb = COLOR_ACCENT
    g_card.line.width = Pt(2)

    gtf = g_card.text_frame
    gtf.word_wrap = True
    gtf.margin_left = Inches(0.4)
    gtf.margin_top = Inches(0.3)

    gp1 = gtf.paragraphs[0]
    gp1.text = "🛡️ 이어잡 인사담당자 안심 보장제 (Risk-Zero Policy)"
    gp1.font.name = FONT_FAMILY
    gp1.font.size = Pt(16)
    gp1.font.bold = True
    gp1.font.color.rgb = COLOR_ACCENT

    gp2 = gtf.add_paragraph()
    gp2.text = "\n\"채용 후 14일 이내 업무 적합도 불만족 시, 추가 비용 없이 100% 무상으로 1회 재매칭을 보장합니다.\""
    gp2.font.name = FONT_FAMILY
    gp2.font.size = Pt(13)
    gp2.font.bold = True
    gp2.font.color.rgb = COLOR_WHITE

    gp3 = gtf.add_paragraph()
    gp3.text = "이제 시니어 전문가 채용에 따르는 역량 검증 실패 리스크와 비용 부담을 이어잡이 함께 분담해 드립니다."
    gp3.font.name = FONT_FAMILY
    gp3.font.size = Pt(11)
    gp3.font.color.rgb = COLOR_MINT

    # --------------------------------------------------------------------------
    # SLIDE 8: CTA & 특별 혜택 (High-End Dark Closing)
    # --------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_bg_color(s8, COLOR_PRIMARY)

    # Center Modal Container
    m_card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.0), Inches(9.733), Inches(5.5))
    m_card.fill.solid()
    m_card.fill.fore_color.rgb = COLOR_WHITE
    m_card.line.color.rgb = COLOR_ACCENT
    m_card.line.width = Pt(2)

    if os.path.exists(logo_path):
        s8.shapes.add_picture(logo_path, Inches(5.9), Inches(1.4), height=Inches(0.42))

    c_box = s8.shapes.add_textbox(Inches(2.2), Inches(2.0), Inches(8.933), Inches(2.2))
    c_tf = c_box.text_frame
    c_tf.word_wrap = True

    cp1 = c_tf.paragraphs[0]
    cp1.text = "지금 우리 회사에 꼭 필요한\n검증된 시니어 인재 3인을 무료로 추천받으세요"
    cp1.font.name = FONT_FAMILY
    cp1.font.size = Pt(23)
    cp1.font.bold = True
    cp1.font.color.rgb = COLOR_PRIMARY
    cp1.alignment = PP_ALIGN.CENTER

    cp2 = c_tf.add_paragraph()
    cp2.text = "\n채용 중인 포지션 조건이나 당면 과제를 메일로 회신해 주시면,\n최적 매칭 인재 3인의 AI 경험 카드와 고용촉진장려금 산출서를 24시간 내 무상으로 전달해 드립니다."
    cp2.font.name = FONT_FAMILY
    cp2.font.size = Pt(12.5)
    cp2.font.color.rgb = COLOR_INK_MUTED
    cp2.alignment = PP_ALIGN.CENTER

    # Action Button Pill in Modal
    act_pill = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.55))
    act_pill.fill.solid()
    act_pill.fill.fore_color.rgb = COLOR_ACCENT
    act_pill.line.fill.background()
    act_p = act_pill.text_frame.paragraphs[0]
    act_p.text = "📩 이메일 회신으로 무료 신청하기"
    act_p.font.name = FONT_FAMILY
    act_p.font.size = Pt(12)
    act_p.font.bold = True
    act_p.font.color.rgb = COLOR_WHITE
    act_p.alignment = PP_ALIGN.CENTER

    # Contact Info Bar in Modal
    c_info_box = s8.shapes.add_textbox(Inches(2.2), Inches(5.15), Inches(8.933), Inches(0.8))
    c_info_tf = c_info_box.text_frame
    cip = c_info_tf.paragraphs[0]
    cip.text = "🌐 공식 웹사이트: https://eojob.kr     |     ✉️ 제휴/채용 문의: contact@eojob.kr"
    cip.font.name = FONT_FAMILY
    cip.font.size = Pt(11.5)
    cip.font.bold = True
    cip.font.color.rgb = COLOR_PRIMARY
    cip.alignment = PP_ALIGN.CENTER

    # Enforce Pretendard Everywhere
    enforce_pretendard_everywhere(prs)

    output_path = os.path.abspath("docs/slides/eojob-intro.pptx")
    prs.save(output_path)
    print(f"High-end presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_high_end_presentation()
